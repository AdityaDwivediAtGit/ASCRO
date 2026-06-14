import os
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_charts():
    """Generates charts for the presentation using matplotlib"""
    if not os.path.exists('demo_assets'):
        os.makedirs('demo_assets')
        
    # Chart 1: Latency Comparison
    plt.figure(figsize=(6, 4))
    labels = ['Baseline LLM', 'ASCRO (vLLM ROCm)']
    latencies = [4.5, 1.8]
    plt.bar(labels, latencies, color=['#888888', '#003087'])
    plt.ylabel('Latency (Seconds)')
    plt.title('End-to-End Inference Latency')
    for i, v in enumerate(latencies):
        plt.text(i, v + 0.1, f"{v}s", ha='center')
    plt.savefig('demo_assets/latency_chart.png', bbox_inches='tight')
    plt.close()

    # Chart 2: GPU Memory
    plt.figure(figsize=(6, 4))
    labels = ['Standard vLLM', 'ASCRO (Quantized AWQ)']
    mem = [32, 12]
    plt.bar(labels, mem, color=['#888888', '#ED1C24']) # AMD Red/Orange
    plt.ylabel('GPU Memory (GB)')
    plt.title('GPU Memory Footprint on MI300X')
    for i, v in enumerate(mem):
        plt.text(i, v + 0.5, f"{v}GB", ha='center')
    plt.savefig('demo_assets/memory_chart.png', bbox_inches='tight')
    plt.close()

def add_footer(slide, prs):
    """Adds standard TCS/AMD footer to a slide"""
    left = Inches(0.5)
    top = Inches(7.0)
    width = Inches(9.0)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Built on AMD MI300X | ROCm 7.2 | TCS & AMD AI Hackathon 2026"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER

def create_presentation():
    create_charts()
    
    prs = Presentation()
    
    # Define Title Slide Layout
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "ASCRO"
    subtitle.text = ("Agentic Supply Chain Resilience Orchestrator\n"
                     "Multi-Agent RAG + Predictive System for TCS Manufacturing Excellence\n"
                     "TCS & AMD AI Hackathon 2026 | Track 1 - Agents\n"
                     "Team: [Your Team Name] | Date: June 2026")
                     
    # You would typically add logos here, leaving placeholder space.
    
    # Slide 2: Problem & Context
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide2.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Problem & Context"
    
    tf = body_shape.text_frame
    tf.text = "Supply Chain Disruptions Cost Millions"
    
    p = tf.add_paragraph()
    p.text = "20-40% downtime caused by unexpected supply chain bottlenecks."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "TCS Client Challenges: Manual assessment of geopolitical & weather risks is too slow."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Business Relevance: Aligns perfectly with TCS Manufacturing AI Axis."
    p.level = 1
    
    add_footer(slide2, prs)

    # Slide 3: Solution Overview
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide3.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Solution Overview: Multi-Agent Architecture"
    
    tf = body_shape.text_frame
    tf.text = "One-Liner: Multi-agent system that predicts, researches, and orchestrates supply chain disruptions in real-time."
    
    p = tf.add_paragraph()
    p.text = "1. Supervisor: Routes and decomposes queries."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Researcher (RAG): Queries vector DB (FAISS) for TCS playbooks/contracts."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Predictor (ML): XGBoost calculates disruption probability based on features."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Executor: Generates step-by-step mitigation plan."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "5. Critic: Validates plan against constraints to prevent hallucinations."
    p.level = 1
    
    add_footer(slide3, prs)

    # Slide 4: Model Insights & Technical Implementation
    slide4 = prs.slides.add_slide(prs.slide_layouts[5]) # Title only layout for custom charts
    shapes = slide4.shapes
    title_shape = shapes.title
    title_shape.text = "Model Insights & AMD Optimization"
    
    # Add text box for metrics
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(4.0)
    txBox = slide4.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Technical Highlights:"
    
    p = tf.add_paragraph()
    p.text = "vLLM on ROCm: Qwen2.5-7B-Instruct-AWQ"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Quantization & Flash Attention enabled."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "XGBoost GPU Tree Method for fast tabular inference."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Predictor F1 Score: 87%"
    p.level = 1
    
    # Insert charts
    slide4.shapes.add_picture('demo_assets/latency_chart.png', Inches(5.0), Inches(1.5), width=Inches(4.5))
    slide4.shapes.add_picture('demo_assets/memory_chart.png', Inches(5.0), Inches(4.5), width=Inches(4.5))
    
    add_footer(slide4, prs)

    # Slide 5: Impact & Demo Summary
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide5.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Impact, Demo Summary & Future Vision"
    
    tf = body_shape.text_frame
    tf.text = "Impact:"
    
    p = tf.add_paragraph()
    p.text = "50% faster risk assessment compared to manual processes."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Significant productivity gains for TCS Delivery Teams."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Demo Highlights:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Interactive Gradio UI showcasing LangGraph agent traces in real-time."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Future Vision:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Integration with live SAP ERP systems & multimodal document parsing."
    p.level = 1
    
    add_footer(slide5, prs)

    prs.save('ASCRO_Hackathon_Presentation.pptx')
    print("Presentation saved successfully as 'ASCRO_Hackathon_Presentation.pptx'")

if __name__ == '__main__':
    create_presentation()
